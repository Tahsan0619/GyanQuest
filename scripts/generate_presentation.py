"""
Generate a PowerPoint presentation for the Force Fighter pitch.

Requires: python-pptx

Usage:
  pip install python-pptx
  python scripts/generate_presentation.py --out Force_Fighter_Pitch.pptx

This script creates 5 slides with speaker notes based on repo prompts.
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle=None, notes=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_bulleted_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_presentation(path):
    prs = Presentation()

    # Slide 1 — Vision
    add_title_slide(
        prs,
        title="From one village classroom to 170 million learners",
        subtitle="The vision: one country, one digital learning layer",
        notes=(
            "Problem + vision — phones yes, labs no.\n"
            "Show Force Fighter — live demo: wake rock + playground + বাংলা toggle."
        ),
    )

    # Slide 2 — Platform
    add_bulleted_slide(
        prs,
        title="Build once, scale to every subject",
        bullets=[
            "Roadmap: 27 domains across Sciences, CS, Engineering, Advanced Science, Math",
            "Force Fighter proves: 10-level mission path, playground, bilingual UI, responsive",
            "Shared shell: i18n, progress, coach UI; per-subject content + simulation",
            "One line: Physics is live proof — next 26 are content modules, not new products",
        ],
        notes=(
            "Comparison table — “why not YouTube?” emphasize playground + local language."
        ),
    )

    # Slide 3 — Competitive advantage
    add_bulleted_slide(
        prs,
        title="We win on what Bangladesh actually needs",
        bullets=[
            "Active missions vs passive video; local context + বাংলা-first examples",
            "Browser-first: no hardware, low-cost scaling, offline after first load",
            "Pedagogy: concrete-before-abstract, safe failure, spiral curriculum",
            "Marginal cost → zero per student once built",
        ],
        notes=(
            "Comparison table — highlight each competitor weakness and our advantage."
        ),
    )

    # Slide 4 — Built for Bangladesh
    add_bulleted_slide(
        prs,
        title="National scale, not elite scale",
        bullets=[
            "Mobile-first, low-bandwidth progressive asset loading",
            "বাংলা removes language barrier; playground helps mixed-pace classrooms",
            "Align levels to NCTB chapters; teacher dashboard next phase",
            "Impact metrics: steps completed, time-on-task, quiz pass rate",
        ],
        notes=(
            "Bangladesh-specific — language, rural, NCTB mapping; ask for pilot schools."
        ),
    )

    # Slide 5 — Ask & roadmap
    add_bulleted_slide(
        prs,
        title="The only team building a national interactive curriculum",
        bullets=[
            "Proven product + repeatable factory: 27 subjects on one engine",
            "12-month rollout example: Q1 physics+chem+bio, Q2 math+ICT, Q3 engineering pilot, Q4 full catalog",
            "What we need: pilot schools, subject advisors, hosting/CDN, MoE/NGO partnerships",
            "Closing: Force Fighter is proof — build Bangladesh's national interactive layer",
        ],
        notes=(
            "Roadmap + ask — pilots and partners. 30s wrap: pilots, advisors, hosting."
        ),
    )

    prs.save(path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="Force_Fighter_Pitch.pptx", help="Output pptx file path")
    args = parser.parse_args()
    build_presentation(args.out)
    print(f"Saved presentation to {args.out}")


if __name__ == "__main__":
    main()
