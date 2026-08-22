#!/usr/bin/env python3
"""
GyanQuest BrainChild 2.0 pitch deck (8 slides).

Source of truth: GyanQuest_Project_Report.pdf (not the older DOCX).
Pedagogy mapping: engine/js/book-theory.js, engine/js/pedagogy.js, persist.js.

Design: useslidekit ANTI-SLOP + DESIGN-SYSTEM, antislop-decks.
  - Assertion titles (complete sentences)
  - Solid fills, no gradients, no title underlines, no drop shadows
  - Sharp rectangles (editorial, 0 radius)
  - Georgia display + Cambria body (not Inter / Roboto / Calibri)
  - Layouts rotate; gold is reserved for the 80% lock
  - No Thank-you slide; demo is the close
  - No em dashes
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "GyanQuest_BrainChild_Pitch.pptx"

# Report palette. Gold is a lock, not a decoration under titles.
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
INK = RGBColor(0x1B, 0x28, 0x38)
MUTED = RGBColor(0x5C, 0x6B, 0x7A)
IVORY = RGBColor(0xF4, 0xEB, 0xD0)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW = RGBColor(0xEE, 0xF3, 0xF8)
BAND = RGBColor(0xE8, 0xE0, 0xD0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 8
DISPLAY = "Georgia"
BODY = "Cambria"


def _rfonts(run, name: str) -> None:
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def _set_run(run, size, bold=False, color=INK, italic=False, font=BODY):
    _rfonts(run, font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.shadow = False


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _box(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, fill)
    sh.shadow.inherit = False
    return sh


def _tb(
    slide,
    l,
    t,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    color=INK,
    font=BODY,
    align=PP_ALIGN.LEFT,
    italic=False,
    anchor="t",
):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {"t": "t", "ctr": "ctr", "b": "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    _set_run(run, size, bold, color, italic, font)
    return box


def _lines(slide, l, t, w, h, items, *, size=16, color=INK, gap=8, font=BODY, bold=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = item
        _set_run(run, size, bold, color, False, font)
    return box


def _footer(slide, n, *, light=True):
    c = MUTED if light else RGBColor(0xA8, 0xB4, 0xC0)
    _tb(
        slide,
        Inches(0.55),
        Inches(7.18),
        Inches(10.5),
        Inches(0.26),
        "BrainChild 2.0  ·  Team Alpha  ·  University of Frontier Technology, Bangladesh",
        size=11,
        color=c,
        font=BODY,
    )
    _tb(
        slide,
        Inches(11.45),
        Inches(7.18),
        Inches(1.3),
        Inches(0.26),
        f"{n} / {TOTAL}",
        size=11,
        color=c,
        align=PP_ALIGN.RIGHT,
        font=BODY,
    )


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rail(slide):
    _box(slide, Inches(0), Inches(0), Inches(0.16), SLIDE_H, NAVY)


def _block(slide, l, t, w, h, lines, *, size=13, color=INK, font=BODY, bold_first=False):
    """Multi-line text box. Each string is its own paragraph (python-pptx ignores \\n in a run)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    parts = []
    for line in lines:
        parts.extend(str(line).split("\n"))
    for i, line in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        _set_run(run, size, bold_first and i == 0, color, False, font)
    return box


def _kicker(slide, text):
    _tb(
        slide,
        Inches(0.55),
        Inches(0.28),
        Inches(12.2),
        Inches(0.28),
        text,
        size=12,
        color=MUTED,
        bold=True,
        font=BODY,
    )


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    prs.core_properties.title = "GyanQuest · BrainChild 2.0"
    prs.core_properties.author = "Team Alpha"
    prs.core_properties.subject = "Pitch deck. Source: GyanQuest_Project_Report.pdf"
    prs.core_properties.category = "BrainChild 2.0"

    # ------------------------------------------------------------------ 1 Bruner
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "JEROME S. BRUNER  ·  TOWARD A THEORY OF INSTRUCTION (1966)")
    _tb(
        s,
        Inches(0.55),
        Inches(0.58),
        Inches(12.2),
        Inches(1.15),
        "Knowledge is built in three modes, then spiraled back richer.",
        size=32,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    _tb(
        s,
        Inches(0.55),
        Inches(1.78),
        Inches(12.2),
        Inches(0.62),
        "A learner can act an idea with the hands, see it as a picture, then name it as a rule. "
        "The spiral revisits the same idea at a higher grain. GyanQuest writes this on every mission intro: do and see, then pictures, then name the rule.",
        size=16,
        color=INK,
    )
    modes = [
        (BAND, "1   Enactive", "Learn by doing.", "Hands on the stage: drag, tap, snap, sort, fail, retry."),
        (ROW, "2   Iconic", "Learn by seeing.", "Pictures and motion: chips, dials, 2D labs, 3D specimen."),
        (BAND, "3   Symbolic", "Learn by naming.", "Language and form: the rule, the formula, the 80% drill."),
    ]
    widths = [12.22, 10.55, 8.85]
    y = 2.52
    for (fill, title, why, runs), w in zip(modes, widths):
        _box(s, Inches(0.55), Inches(y), Inches(w), Inches(1.22), fill)
        _tb(s, Inches(0.75), Inches(y + 0.10), Inches(w - 0.4), Inches(0.34), title, size=20, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(0.75), Inches(y + 0.44), Inches(w - 0.4), Inches(0.26), why, size=14, italic=True, color=MUTED, font=BODY)
        _tb(s, Inches(0.75), Inches(y + 0.72), Inches(w - 0.4), Inches(0.38), runs, size=16, color=INK, font=BODY)
        y += 1.34
    _footer(s, 1)
    _notes(
        s,
        "Open with Bruner, not the product. Three modes: enactive (do), iconic (see), symbolic (name). "
        "Spiral curriculum: same idea, higher grain, not a harder chapter dumped as level 10. "
        "Then: GyanQuest is this spiral made into software. The CSS class in code is chem-intro__brunner; "
        "the citation in the PDF is Bruner 1966.",
    )

    # ------------------------------------------------------------------ 2 Objective
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "OBJECTIVES  ·  FROM THE BRAINCHILD REPORT")
    _tb(
        s,
        Inches(0.55),
        Inches(0.58),
        Inches(12.2),
        Inches(1.05),
        "A school outcome becomes a mission with a win condition.",
        size=30,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    objs = [
        ("01", "Playable, not watched", "Each stated objective has a clear win: sort, snap, drill, mastery."),
        ("02", "Readable stage", "Canvas 2D for graded missions. 3D only where spatial inspection helps."),
        ("03", "State that survives reload", "save-v2 restores hub vs in-mission in the browser."),
        ("04", "Bangla and English first-class", "Two language modes. Not a subtitle overlay."),
        ("05", "AI cannot block class", "If Groq is down or the key is missing, the game still teaches."),
        ("06", "Browser is the lab", "No install. A modern browser is the only required hardware."),
    ]
    for i, (n, h, b) in enumerate(objs):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.4
        yy = 1.82 + row * 1.68
        _tb(s, Inches(x), Inches(yy), Inches(0.7), Inches(0.4), n, size=20, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(x + 0.75), Inches(yy), Inches(5.3), Inches(0.4), h, size=18, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(x + 0.75), Inches(yy + 0.42), Inches(5.3), Inches(1.05), b, size=16, color=INK, font=BODY)
    _footer(s, 2)
    _notes(
        s,
        "Six objectives, PDF section 3. Stress 05: Groq down still means class runs. Stress 02: 3D is Specimen Lab, "
        "not every lesson. Ages 9-16, village computer class and city coaching center run the same titles.",
    )

    # ------------------------------------------------------------------ 3 Problem
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "PROBLEM STATEMENT")
    _tb(
        s,
        Inches(0.55),
        Inches(0.58),
        Inches(7.55),
        Inches(1.45),
        "Lecture-heavy STEAM never proves the idea was enacted.",
        size=30,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    _tb(
        s,
        Inches(0.55),
        Inches(2.15),
        Inches(7.55),
        Inches(1.35),
        "Bangladeshi classrooms: shared phones, flaky internet, no dedicated hardware lab. "
        "Tools are a video here, a quiz there, a 3D demo that only runs on a lab PC.",
        size=17,
        color=INK,
    )
    facts = [
        "Learners watch and click MCQ. Teachers cannot see enacted vs memorized.",
        "A village computer class and a city coaching center cannot share one product.",
        "Time-watched is treated as progress. Fluency is not.",
    ]
    yy = 3.62
    for f in facts:
        _box(s, Inches(0.55), Inches(yy), Inches(0.14), Inches(0.78), GOLD)
        _tb(s, Inches(0.86), Inches(yy + 0.08), Inches(7.15), Inches(0.7), f, size=16, color=INK, font=BODY)
        yy += 0.98
    _box(s, Inches(8.45), Inches(0.28), Inches(4.4), Inches(6.65), NAVY)
    _tb(s, Inches(8.7), Inches(0.55), Inches(3.95), Inches(0.4), "What fails today", size=14, bold=True, color=GOLD, font=BODY)
    fails = [
        ("Watch + MCQ", "not act on a stage"),
        ("Time watched", "not an 80% lock"),
        ("Cloud as a gate", "not a bonus"),
        ("3D everywhere or never", "not a lab beside 2D"),
        ("English UI, Bangla afterthought", "not two modes"),
        ("One app per topic", "not one engine"),
    ]
    yy = 1.15
    for left, right in fails:
        _tb(s, Inches(8.7), Inches(yy), Inches(3.95), Inches(0.32), left, size=16, bold=True, color=IVORY, font=DISPLAY)
        _tb(s, Inches(8.7), Inches(yy + 0.32), Inches(3.95), Inches(0.32), right, size=14, color=RGBColor(0xC5, 0xD0, 0xDC), font=BODY)
        yy += 0.88
    _footer(s, 3)
    _notes(
        s,
        "PDF section 2. Cover banner and problem say STEAM; the abstract still says STEM once. "
        "Say STEAM out loud. Right panel is PDF section 12 (typical EdTech vs GyanQuest), inverted as the failure mode. "
        "Do not attack named products.",
    )

    # ------------------------------------------------------------------ 4 Product
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "THE PRODUCT")
    _tb(
        s,
        Inches(0.55),
        Inches(0.55),
        Inches(12.2),
        Inches(0.85),
        "GyanQuest is one engine, 28 live titles, four layers.",
        size=30,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    stats = [
        ("28", "subject titles"),
        ("16", "uniqueness-pass"),
        ("25", "finished missions"),
        ("80%", "fluency gate"),
        ("2", "language modes"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.55 + i * 2.52
        _tb(s, Inches(x), Inches(1.5), Inches(2.35), Inches(0.7), num, size=36, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(x), Inches(2.2), Inches(2.35), Inches(0.35), label, size=13, color=MUTED, font=BODY)
    _tb(
        s,
        Inches(0.55),
        Inches(2.65),
        Inches(12.2),
        Inches(0.35),
        "Landing  →  mission hub  →  Canvas 2D  →  coach  →  digital book  →  optional Groq tutor. Progress stays local.",
        size=15,
        color=INK,
        font=BODY,
    )
    tracks = [
        ("Core science  5", "Force, Chemistry, Bio, Math Quest, Eco"),
        ("CS and tech  10", "ICT, Web, Backend, SQL, Net, Cyber, OS, AI, ML, Data"),
        ("Engineering  5", "Electrical, Mechanical, Civil, Robotics, Green Tech"),
        ("Adv. science  4", "Astronomy, Geology, Anatomy, Genetics"),
        ("Math extended  4", "Statistics, Geometry, Calculus, Discrete"),
    ]
    for i, (h, b) in enumerate(tracks):
        yy = 3.12 + i * 0.48
        _tb(s, Inches(0.55), Inches(yy), Inches(2.9), Inches(0.42), h, size=14, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(3.5), Inches(yy), Inches(4.6), Inches(0.42), b, size=14, color=INK, font=BODY)
    layers = [
        ("Canvas", "Source of truth. Drag, tap, myth-bust."),
        ("Coach", "Streaks and misses. Live commentary."),
        ("Book", "Retrieval + dual coding. Term links."),
        ("AI tutor", "On demand. Proxy + fallback. Never blocking."),
    ]
    _box(s, Inches(8.4), Inches(3.12), Inches(4.4), Inches(3.7), NAVY)
    _tb(s, Inches(8.62), Inches(3.28), Inches(4.0), Inches(0.32), "Four layers", size=13, bold=True, color=GOLD, font=BODY)
    for i, (h, b) in enumerate(layers):
        yy = 3.72 + i * 0.74
        _tb(s, Inches(8.62), Inches(yy), Inches(4.0), Inches(0.3), f"{i + 1}  {h}", size=16, bold=True, color=IVORY, font=DISPLAY)
        _tb(s, Inches(8.62), Inches(yy + 0.30), Inches(4.0), Inches(0.38), b, size=13, color=RGBColor(0xC5, 0xD0, 0xDC), font=BODY)
    _footer(s, 4)
    _notes(
        s,
        "August 2026 catalog from the PDF table. Do not read all 28 names. Honesty from the report: not every "
        "28×10×10 cell is a unique authored story yet. Force Fighter and uniqueness-pass titles carry the deepest labs. "
        "Stubs exist on purpose. Four layers: canvas leads; AI is optional.",
    )

    # ------------------------------------------------------------------ 5 STEAM / solving
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "WHAT IT SOLVES  ·  STEAM, NOT STEM-ONLY")
    _tb(
        s,
        Inches(0.55),
        Inches(0.55),
        Inches(12.2),
        Inches(0.95),
        "It solves STEAM enactment in a browser, not another STEM quiz.",
        size=28,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    steam = [
        ("S", NAVY, IVORY, "Science", "Chem, bio, astronomy, geology, anatomy, genetics: act the model, then name it."),
        ("T", NAVY, IVORY, "Technology", "ICT, web, SQL, cyber, AI, ML: one hub, not eight separate apps."),
        ("E", NAVY, IVORY, "Engineering", "Electrical, mechanical, civil, robotics, green tech: loops, levers, loads."),
        ("A", GOLD, NAVY, "Arts", "Bilingual craft, digital books, myth stories, Specimen Lab as a visual theater."),
        ("M", NAVY, IVORY, "Mathematics", "Number, stats, geometry, calculus, discrete: drill to 80% before mastery."),
    ]
    yy = 1.62
    for letter, fill, tc, name, line in steam:
        _box(s, Inches(0.55), Inches(yy), Inches(0.78), Inches(0.82), fill)
        _tb(
            s,
            Inches(0.55),
            Inches(yy + 0.16),
            Inches(0.78),
            Inches(0.52),
            letter,
            size=26,
            bold=True,
            color=tc,
            align=PP_ALIGN.CENTER,
            font=DISPLAY,
        )
        _tb(s, Inches(1.52), Inches(yy + 0.04), Inches(2.4), Inches(0.32), name, size=18, bold=True, color=NAVY, font=DISPLAY)
        _tb(s, Inches(3.95), Inches(yy + 0.12), Inches(8.8), Inches(0.62), line, size=16, color=INK, font=BODY)
        yy += 0.92
    _tb(
        s,
        Inches(0.55),
        Inches(6.35),
        Inches(12.2),
        Inches(0.55),
        "Also solved: offline-first access, 80% fluency vs time-watched, 2D missions plus a 3D lab instead of 3D everywhere.",
        size=15,
        italic=True,
        color=MUTED,
        font=BODY,
    )
    _footer(s, 5)
    _notes(
        s,
        "User and cover: STEAM. Arts is not a drawing class. Arts in this product: bilingual writing, book craft, "
        "myth narrative, dual-coded visuals, Specimen Lab as inspection theater. "
        "Impact line from PDF: converting curriculum into mastery-gated missions at the cost of a browser, not a hardware lab. "
        "Pause on A. Then fluency lock.",
    )

    # ------------------------------------------------------------------ 6 Theories
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "THEORIES IN THE RUNNING PRODUCT")
    _tb(
        s,
        Inches(0.55),
        Inches(0.52),
        Inches(12.2),
        Inches(0.7),
        "Each theory maps to a part that actually runs.",
        size=28,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    headers = ("Theory", "Why it is here", "Component", "What actually runs")
    col_x = [0.45, 3.15, 6.55, 9.35]
    col_w = [2.62, 3.28, 2.68, 3.5]
    _box(s, Inches(0.45), Inches(1.28), Inches(12.4), Inches(0.38), NAVY)
    for x, w, hdg in zip(col_x, col_w, headers):
        _tb(s, Inches(x + 0.08), Inches(1.32), Inches(w - 0.1), Inches(0.3), hdg, size=12, bold=True, color=WHITE, font=BODY)
    rows = [
        (
            "Bruner 1966\nspiral + 3 modes",
            "Move from doing to naming without skipping the picture.",
            "Mission intro +\n10-step spiral",
            "Enactive labs → iconic chips/3D → symbolic rule + drill",
        ),
        (
            "Constructivism +\nPOE (White & Gunstone)",
            "Canvas must be the source of truth, not a video.",
            "Interactive 2D stage",
            "Predict, act, observe, explain; drag / tap / myth-bust",
        ),
        (
            "Vygotsky ZPD +\nFlow + formative",
            "Help at the edge of skill, not a lecture dump.",
            "Game coach",
            "3-hit streak praise; 2-miss simpler take (streak-based)",
        ),
        (
            "Paivio, Mayer,\nSweller CLT",
            "Words and pictures together; one idea per page.",
            "Digital mission book",
            "8-page spine: hook → model → myth → retrieval",
        ),
        (
            "Posner conceptual\nchange",
            "A misconception must be faced, not hidden.",
            "Myth-bust step",
            "Claim vs better model in step 8 of the spiral",
        ),
        (
            "Retrieval +\nBloom 80%",
            "Mastery is earned. Time watched is not progress.",
            "Fluency drill lock",
            "Step 10 stays closed until drill hits 80%",
        ),
        (
            "Cognitive\napprenticeship",
            "Tutor on demand, never the lesson itself.",
            "Groq proxy +\nlocal fallback",
            "Socratic chat; canvas and book remain the truth",
        ),
    ]
    yy = 1.68
    rh = 0.66
    for i, row in enumerate(rows):
        _box(s, Inches(0.45), Inches(yy), Inches(12.4), Inches(rh), ROW if i % 2 else WHITE)
        for j, (x, w, val) in enumerate(zip(col_x, col_w, row)):
            _block(
                s,
                Inches(x + 0.08),
                Inches(yy + 0.05),
                Inches(w - 0.12),
                Inches(rh - 0.08),
                [val],
                size=13,
                color=INK,
                font=BODY,
                bold_first=(j == 0),
            )
        yy += rh
    _tb(
        s,
        Inches(0.55),
        Inches(6.38),
        Inches(12.2),
        Inches(0.5),
        "Gee (2007): the game is the literacy, not a wrapper around a quiz. Honest limit: Flow Autopilot is streak-based, not a full adaptive engine.",
        size=13,
        italic=True,
        color=MUTED,
        font=BODY,
    )
    _footer(s, 6)
    _notes(
        s,
        "Walk left to right on Bruner if time is short. Sources: PDF layers table + engine/js/book-theory.js + pedagogy.js. "
        "Coach: 3 correct = on a roll; 2 wrong = simpler take. Mastery locked in persist.js canEnterMastery at 0.8. "
        "Do not invent extra theories. Do not claim a full adaptive engine.",
    )

    # ------------------------------------------------------------------ 7 Learning flow (short)
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rail(s)
    _kicker(s, "LEARNING FLOW")
    _tb(
        s,
        Inches(0.55),
        Inches(0.55),
        Inches(12.2),
        Inches(0.85),
        "Bruner’s spiral sits inside a locked cycle.",
        size=32,
        bold=True,
        color=NAVY,
        font=DISPLAY,
    )
    steps = [
        ("Recall", NAVY, IVORY),
        ("Predict", NAVY, IVORY),
        ("Act", NAVY, IVORY),
        ("Feedback", NAVY, IVORY),
        ("80% fluency", GOLD, NAVY),
        ("Mastery", NAVY, IVORY),
    ]
    for i, (label, fill, tc) in enumerate(steps):
        x = 0.55 + i * 2.12
        _box(s, Inches(x), Inches(1.6), Inches(1.95), Inches(1.15), fill)
        _tb(
            s,
            Inches(x),
            Inches(1.92),
            Inches(1.95),
            Inches(0.55),
            label,
            size=16,
            bold=True,
            color=tc,
            align=PP_ALIGN.CENTER,
            font=DISPLAY,
            anchor="ctr",
        )
    _tb(
        s,
        Inches(0.55),
        Inches(3.0),
        Inches(12.2),
        Inches(0.4),
        "Ten steps inside each topic (Force Fighter pattern)",
        size=15,
        italic=True,
        color=MUTED,
        font=BODY,
    )
    _tb(
        s,
        Inches(0.55),
        Inches(3.4),
        Inches(12.2),
        Inches(0.7),
        "Hook  ·  Watch  ·  Sort  ·  Try  ·  Explain  ·  Name the rule  ·  Stretch  ·  Myth-bust  ·  Fluency  ·  Mastery",
        size=16,
        color=INK,
        font=BODY,
    )
    _box(s, Inches(0.55), Inches(4.3), Inches(12.22), Inches(2.5), ROW)
    _tb(
        s,
        Inches(0.8),
        Inches(4.5),
        Inches(11.7),
        Inches(2.1),
        "Bruner in that list: Try is enactive. Watch, chips, and Specimen Lab are iconic. "
        "Name the rule and the fluency drill are symbolic. Mastery stays locked until the drill is earned. "
        "That is a product rule, not a slogan.",
        size=18,
        color=INK,
        font=BODY,
    )
    _footer(s, 7)
    _notes(
        s,
        "Keep this short. Point at the gold 80% box. Engine cycle from PDF section 6 and pedagogy.js runPreMission: "
        "recall → objective compass → predict → act → feedback → fluency → mastery. Then: we will show it live.",
    )

    # ------------------------------------------------------------------ 8 Demo
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _box(s, Inches(0), Inches(0), SLIDE_W, Inches(0.14), GOLD)
    _tb(
        s,
        Inches(0.7),
        Inches(1.85),
        Inches(12.0),
        Inches(0.35),
        "BRAINCHILD 2.0  ·  TEAM ALPHA",
        size=14,
        bold=True,
        color=GOLD,
        font=BODY,
    )
    _tb(
        s,
        Inches(0.7),
        Inches(2.25),
        Inches(12.0),
        Inches(1.15),
        "Demo starts here.",
        size=48,
        bold=True,
        color=IVORY,
        font=DISPLAY,
    )
    _tb(
        s,
        Inches(0.7),
        Inches(3.55),
        Inches(12.0),
        Inches(1.35),
        "Play one 2D mission through the spiral. Open Specimen Lab and tap a pin. Switch English / Bangla.",
        size=22,
        color=WHITE,
        font=BODY,
    )
    _tb(
        s,
        Inches(0.7),
        Inches(5.15),
        Inches(12.0),
        Inches(0.7),
        "Canvas and book stay the source of truth. The tutor is optional.",
        size=16,
        italic=True,
        color=RGBColor(0xA8, 0xB4, 0xC0),
        font=BODY,
    )
    _footer(s, 8, light=False)
    _notes(
        s,
        "Do not thank the jury on this slide. Open the browser. Landing, one Force Fighter or Bio Explorer mission "
        "through fluency if time, Specimen Lab pin on animal cell, language toggle. "
        "PDF conclusion: we ask the jury to play a mission, pin a cell, and switch language. That is the product.",
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
